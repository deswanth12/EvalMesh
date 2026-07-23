"""
EvalMesh Investor Pitch Deck Generator (Silicon Valley VC Grade)
Creates an ultra-professional, 10-slide PowerPoint (.pptx) presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen format (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Institutional VC Palette
    COLOR_BG = RGBColor(11, 15, 25)          # Deep Slate Dark #0B0F19
    COLOR_CARD = RGBColor(17, 24, 39)        # Translucent Slate #111827
    COLOR_CARD_BORDER = RGBColor(31, 41, 55) # Border Gray
    COLOR_CYAN = RGBColor(0, 240, 255)       # Electric Cyan #00F0FF
    COLOR_BLUE = RGBColor(59, 130, 246)      # Accent Blue #3B82F6
    COLOR_PURPLE = RGBColor(139, 92, 246)   # Accent Purple #8B5CF6
    COLOR_GREEN = RGBColor(16, 185, 129)    # Neon Emerald #10B981
    COLOR_RED = RGBColor(239, 68, 68)       # Crimson Warning #EF4444
    COLOR_TEXT_MAIN = RGBColor(249, 250, 251) # Clean White #F9FAFB
    COLOR_TEXT_MUTED = RGBColor(156, 163, 175) # Muted Silver #9CA3AF

    def set_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_header(slide, title_text, tag_text="EVALMESH | AI AGENT INFRASTRUCTURE"):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = tag_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_CYAN

        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.75))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_MAIN

    # ==========================================
    # SLIDE 1: Cover / Hero Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_bg(slide1)

    # Hero Box Outer Frame
    hero = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.2), Inches(10.93), Inches(5.1))
    hero.fill.solid()
    hero.fill.fore_color.rgb = COLOR_CARD
    hero.line.color.rgb = COLOR_BLUE
    hero.line.width = Pt(2)

    tb1 = slide1.shapes.add_textbox(Inches(1.8), Inches(1.8), Inches(9.7), Inches(3.9))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p1_tag = tf1.paragraphs[0]
    p1_tag.text = "SERIES SEED PITCH DECK"
    p1_tag.font.size = Pt(11)
    p1_tag.font.bold = True
    p1_tag.font.color.rgb = COLOR_CYAN

    p1_main = tf1.add_paragraph()
    p1_main.text = "EvalMesh"
    p1_main.font.size = Pt(52)
    p1_main.font.bold = True
    p1_main.font.color.rgb = COLOR_TEXT_MAIN

    p1_sub = tf1.add_paragraph()
    p1_sub.text = "AI Gateway for Secure & Reliable Agent Deployment"
    p1_sub.font.size = Pt(24)
    p1_sub.font.bold = True
    p1_sub.font.color.rgb = COLOR_BLUE

    p1_desc = tf1.add_paragraph()
    p1_desc.text = "\nCloudflare + GitHub Actions for Autonomous AI Agents"
    p1_desc.font.size = Pt(16)
    p1_desc.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 2: The Problem (Market Pain Points)
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_bg(slide2)
    add_header(slide2, "The Problem: Autonomous Agents Break in Production", "01 / PROBLEM DEFINITION")

    pains = [
        ("💸 Runaway API Billing Spikes", "$2,000+ Surprise Bills", "Recursive agent loops execute 500+ calls overnight without safety limits."),
        ("🔓 Security & PII Data Leaks", "Zero PII Redaction Egress", "Prompt injection jailbreaks trick bots into leaking system secrets & SSNs."),
        ("💥 Silent Model Output Drift", "Broken Downstream Schema", "Upstream model updates alter JSON outputs, crashing client applications.")
    ]

    for i, (title, highlight, desc) in enumerate(pains):
        x = Inches(0.8 + i * 3.9)
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.6), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = COLOR_RED
        card.line.width = Pt(1.5)

        tb = slide2.shapes.add_textbox(x + Inches(0.2), Inches(2.1), Inches(3.2), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(18)
        pt.font.bold = True
        pt.font.color.rgb = COLOR_RED

        ph = tf.add_paragraph()
        ph.text = f"\n{highlight}"
        ph.font.size = Pt(15)
        ph.font.bold = True
        ph.font.color.rgb = COLOR_TEXT_MAIN

        pd = tf.add_paragraph()
        pd.text = f"\n{desc}"
        pd.font.size = Pt(13)
        pd.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 3: The Solution (EvalMesh Gateway)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_bg(slide3)
    add_header(slide3, "The Solution: EvalMesh Zero-Trust Gateway", "02 / THE SOLUTION")

    sols = [
        ("🛡️ Cloudflare for AI", "<15ms Proxy Guardrail", "Inline proxy enforcing real-time WAF, PII DLP redactor, and tool RBAC."),
        ("⚡ Cost Circuit Breaker", "60-90% Cost Reduction", "Semantic cache (3ms, $0) + automatic circuit breaker killing loops at depth > 25."),
        ("🔄 GitHub Actions for AI", "Continuous Evaluation", "Detects semantic output drift & auto-corrects malformed JSON schema outputs.")
    ]

    for i, (title, highlight, desc) in enumerate(sols):
        x = Inches(0.8 + i * 3.9)
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.6), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = COLOR_GREEN
        card.line.width = Pt(1.5)

        tb = slide3.shapes.add_textbox(x + Inches(0.2), Inches(2.1), Inches(3.2), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = title
        pt.font.size = Pt(18)
        pt.font.bold = True
        pt.font.color.rgb = COLOR_GREEN

        ph = tf.add_paragraph()
        ph.text = f"\n{highlight}"
        ph.font.size = Pt(15)
        ph.font.bold = True
        ph.font.color.rgb = COLOR_TEXT_MAIN

        pd = tf.add_paragraph()
        pd.text = f"\n{desc}"
        pd.font.size = Pt(13)
        pd.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 4: Architecture & Data Flow
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_bg(slide4)
    add_header(slide4, "Architecture: In-Line Sidecar Proxy Pipeline", "03 / SYSTEM ARCHITECTURE")

    arch_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.73), Inches(4.8))
    arch_box.fill.solid()
    arch_box.fill.fore_color.rgb = COLOR_CARD
    arch_box.line.color.rgb = COLOR_BLUE

    tb4 = slide4.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(11.1), Inches(4.4))
    tf4 = tb4.text_frame
    tf4.word_wrap = True

    p4_1 = tf4.paragraphs[0]
    p4_1.text = "CLIENT APPLICATION (Python SDK / TypeScript SDK / cURL REST)"
    p4_1.font.size = Pt(16)
    p4_1.font.bold = True
    p4_1.font.color.rgb = COLOR_CYAN

    p4_2 = tf4.add_paragraph()
    p4_2.text = "                       │\n                       ▼\n[ 1. Inbound Security Pipeline ]  ──► PII DLP Redactor + Prompt WAF + Tool RBAC\n                       │\n                       ▼\n[ 2. Proxy Execution Engine ]   ──► Semantic Cache (3ms, $0) + Smart Cost Router\n                       │\n                       ▼\n[ 3. Upstream Provider Egress ] ──► OpenAI / Anthropic / DeepSeek HA Failover\n                       │\n                       ▼\n[ 4. Observability & Evals ]   ──► OpenTelemetry Spans + Golden Datasets + Control Panel UI"
    p4_2.font.size = Pt(14)
    p4_2.font.name = "JetBrains Mono"
    p4_2.font.color.rgb = COLOR_TEXT_MAIN

    # ==========================================
    # SLIDE 5: Core Product Capabilities
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_bg(slide5)
    add_header(slide5, "16-Module Core Engine Capability Matrix", "04 / PRODUCT CAPABILITIES")

    features = [
        ("PII DLP Redactor", "Redacts Emails, SSNs, Credit Cards, IPs inline"),
        ("Prompt Injection WAF", "Blocks jailbreak signatures & system overrides"),
        ("Tool RBAC Enforcer", "Restricts API tool execution by agent role"),
        ("Cost Circuit Breaker", "Terminates runaway loops at depth > 25"),
        ("Semantic Cache", "Serves 80%+ similar prompts in <5ms @ $0 cost"),
        ("Smart Cost Router", "Downgrades simple prompts to 15x cheaper GPT-4o-mini"),
        ("HA Provider Failover", "Auto-routes to Anthropic during OpenAI outages"),
        ("Auto-Healing Retries", "Self-corrects malformed JSON outputs")
    ]

    for i, (feat, desc) in enumerate(features):
        row = i // 2
        col = i % 2
        x = Inches(0.8 + col * 5.9)
        y = Inches(1.8 + row * 1.25)

        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(1.1))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = COLOR_BLUE

        tb = slide5.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1), Inches(5.3), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = f"✅ {feat}"
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = COLOR_TEXT_MAIN

        pd = tf.add_paragraph()
        pd.text = desc
        pd.font.size = Pt(12)
        pd.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 6: Enterprise Compliance
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_bg(slide6)
    add_header(slide6, "Enterprise Compliance & Governance Readiness", "05 / ENTERPRISE GOVERNANCE")

    ents = [
        ("🛡️ SOC 2 Type II", "SHA-256 tamper-proof audit trail log exporter"),
        ("🔒 GDPR Compliance", "In-line data minimization & right-to-be-forgotten cleaner"),
        ("🏥 HIPAA BAA Ready", "Scrubs Medical Record Numbers (MRN) & health identifiers"),
        ("🔑 SAML 2.0 / SSO", "Validates Okta & Auth0 enterprise identity tokens"),
        ("☸️ Kubernetes Gateway", "HPA autoscaling deployment spec (k8s-deployment.yaml)")
    ]

    for i, (title, desc) in enumerate(ents):
        y = Inches(1.8 + i * 1.0)
        card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.73), Inches(0.85))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = COLOR_PURPLE

        tb = slide6.shapes.add_textbox(Inches(1.0), y + Inches(0.1), Inches(11.3), Inches(0.65))
        tf = tb.text_frame
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = f"{title}:  {desc}"
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = COLOR_TEXT_MAIN

    # ==========================================
    # SLIDE 7: Market Opportunity (TAM)
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_bg(slide7)
    add_header(slide7, "Market Opportunity: $50B AI Infrastructure Frontier", "06 / MARKET OPPORTUNITY")

    tams = [
        ("$50 Billion", "Total Addressable Market (TAM)", "Enterprise AI security, governance, and infrastructure proxy sidecars by 2028."),
        ("400% YoY", "AI Agent Adoption Growth", "Over 80% of Fortune 500 enterprises are deploying autonomous AI agents in 2026."),
        ("60-90% ROI", "Direct API Dollar Savings", "EvalMesh pays for itself instantly via semantic prompt caching & cost routing.")
    ]

    for i, (stat, title, desc) in enumerate(tams):
        x = Inches(0.8 + i * 3.9)
        card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.6), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = COLOR_CYAN
        card.line.width = Pt(1.5)

        tb = slide7.shapes.add_textbox(x + Inches(0.2), Inches(2.1), Inches(3.2), Inches(4.2))
        tf = tb.text_frame
        tf.word_wrap = True

        ps = tf.paragraphs[0]
        ps.text = stat
        ps.font.size = Pt(36)
        ps.font.bold = True
        ps.font.color.rgb = COLOR_CYAN

        pt = tf.add_paragraph()
        pt.text = f"\n{title}"
        pt.font.size = Pt(16)
        pt.font.bold = True
        pt.font.color.rgb = COLOR_TEXT_MAIN

        pd = tf.add_paragraph()
        pd.text = f"\n{desc}"
        pd.font.size = Pt(13)
        pd.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 8: Business Model & Pricing
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_bg(slide8)
    add_header(slide8, "Business Model & Tiered Monetization", "07 / MONETIZATION MODEL")

    plans = [
        ("Starter", "$0 / mo", "100k req/mo\nPII Redactor\nPrompt WAF\nWeb Dashboard", COLOR_BLUE),
        ("Pro", "$49 / mo", "1M req/mo\nSemantic Cache\nSmart Cost Router\nAuto-Healing", COLOR_GREEN),
        ("Team", "$299 / mo", "10M req/mo\nOTel Exporter\nMulti-Model Router\nPriority Support", COLOR_PURPLE),
        ("Enterprise", "Custom", "Unlimited req\nDedicated VPC\nSLA Guarantee\n24/7 Support", COLOR_CYAN)
    ]

    for i, (name, price, text, color) in enumerate(plans):
        x = Inches(0.8 + i * 2.95)
        card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(2.75), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = color
        card.line.width = Pt(2)

        tb = slide8.shapes.add_textbox(x + Inches(0.15), Inches(2.0), Inches(2.45), Inches(4.4))
        tf = tb.text_frame
        tf.word_wrap = True

        pt = tf.paragraphs[0]
        pt.text = name
        pt.font.size = Pt(22)
        pt.font.bold = True
        pt.font.color.rgb = color

        pp = tf.add_paragraph()
        pp.text = price
        pp.font.size = Pt(26)
        pp.font.bold = True
        pp.font.color.rgb = COLOR_TEXT_MAIN

        pd = tf.add_paragraph()
        pd.text = f"\n{text}"
        pd.font.size = Pt(13)
        pd.font.color.rgb = COLOR_TEXT_MUTED

    # ==========================================
    # SLIDE 9: Traction & Code Verification
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_bg(slide9)
    add_header(slide9, "Production Readiness: 100% Operational Codebase", "08 / TRACTION & VERIFICATION")

    card9 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.73), Inches(4.8))
    card9.fill.solid()
    card9.fill.fore_color.rgb = COLOR_CARD
    card9.line.color.rgb = COLOR_GREEN

    tb9 = slide9.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(11.1), Inches(4.4))
    tf9 = tb9.text_frame
    tf9.word_wrap = True

    p9_1 = tf9.paragraphs[0]
    p9_1.text = "SYSTEM-WIDE AUTOMATED TEST SUITE: 16/16 CHECKS PASSED (100%)"
    p9_1.font.size = Pt(16)
    p9_1.font.bold = True
    p9_1.font.color.rgb = COLOR_GREEN

    p9_2 = tf9.add_paragraph()
    p9_2.text = "\n • Live Investor Demo Script : python live_demo.py\n • 16-Module Audit Suite     : python -m evalmesh.verify_all\n • Open-Source GitHub Repo   : https://github.com/deswanth12/EvalMesh\n • Interactive Web Dashboard : http://localhost:8000\n • Docker & k8s Deployment   : Dockerfile & k8s-deployment.yaml"
    p9_2.font.size = Pt(15)
    p9_2.font.color.rgb = COLOR_TEXT_MAIN

    # ==========================================
    # SLIDE 10: Call to Action & Vision
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_bg(slide10)

    card10 = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.2), Inches(10.93), Inches(5.1))
    card10.fill.solid()
    card10.fill.fore_color.rgb = COLOR_CARD
    card10.line.color.rgb = COLOR_CYAN
    card10.line.width = Pt(2)

    tb10 = slide10.shapes.add_textbox(Inches(1.8), Inches(1.8), Inches(9.7), Inches(3.9))
    tf10 = tb10.text_frame
    tf10.word_wrap = True

    p10_1 = tf10.paragraphs[0]
    p10_1.text = "Join the Future of AI Reliability"
    p10_1.font.size = Pt(42)
    p10_1.font.bold = True
    p10_1.font.color.rgb = COLOR_TEXT_MAIN

    p10_2 = tf10.add_paragraph()
    p10_2.text = "EvalMesh — AI Gateway for Secure & Reliable Agent Deployment"
    p10_2.font.size = Pt(22)
    p10_2.font.bold = True
    p10_2.font.color.rgb = COLOR_CYAN

    p10_3 = tf10.add_paragraph()
    p10_3.text = "\nGitHub Repository : https://github.com/deswanth12/EvalMesh\nLive Dashboard    : http://localhost:8000\nContact Email     : founder@evalmesh.io"
    p10_3.font.size = Pt(16)
    p10_3.font.color.rgb = COLOR_TEXT_MUTED

    # Save Presentation
    output_filename = "EvalMesh_Investor_Pitch_Deck.pptx"
    prs.save(output_filename)
    print(f"[SUCCESS] Successfully generated VC-Grade PowerPoint pitch deck: {output_filename}")

if __name__ == "__main__":
    create_presentation()
